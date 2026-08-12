"""Build the SLEEP presentation deck.

Authors each slide as fixed-canvas HTML (1920x1080), renders each to PNG via
Edge headless (2x scale for crispness), and assembles an image-based PPTX with
full talking points in the speaker notes.

Run:  python build_deck.py
Out:  SLEEP_Presentation.pptx
"""

from __future__ import annotations

import os
import subprocess

HERE = os.path.dirname(os.path.abspath(__file__))
HTML_DIR = os.path.join(HERE, "html")
PNG_DIR = os.path.join(HERE, "png")
FIG = os.path.join(HERE, "..", "paper", "figures")
EDGE = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"

os.makedirs(HTML_DIR, exist_ok=True)
os.makedirs(PNG_DIR, exist_ok=True)


def fig_uri(name: str) -> str:
    p = os.path.abspath(os.path.join(FIG, name)).replace("\\", "/")
    return f"file:///{p}"


CSS = """
* { margin:0; padding:0; box-sizing:border-box; }
body { background:#F4F6FA; }
.slide {
  position:relative; width:1920px; height:1080px; overflow:hidden;
  background:#F4F6FA; color:#1A1F2E;
  font-family:'Segoe UI', system-ui, sans-serif;
}
.slide.dark { background:#141927; color:#E6EAF4; }

.frame { position:absolute; inset:70px 90px 84px 90px; display:flex; flex-direction:column; }
.eyebrow {
  font-family:Consolas, monospace; font-size:24px; letter-spacing:0.16em;
  text-transform:uppercase; color:#3D52C9; margin-bottom:18px;
}
.dark .eyebrow { color:#8FA0FF; }
.eyebrow .era-d { color:#A96D12; }
.eyebrow .era-r { color:#2C7D57; }
h1.title {
  font-family:Georgia, serif; font-weight:normal; font-size:66px;
  line-height:1.12; letter-spacing:-0.01em; max-width:1500px; margin-bottom:44px;
}
h1.title b { font-weight:600; }
.content { flex:1; display:flex; flex-direction:column; gap:30px; min-height:0; }
p.lead { font-size:33px; line-height:1.5; color:#3D4658; max-width:1420px; }
.dark p.lead { color:#C3CAD9; }
p.lead b, li b { font-weight:650; color:#1A1F2E; }
.dark p.lead b { color:#E6EAF4; }

.footer {
  position:absolute; left:90px; right:90px; bottom:34px; display:flex;
  justify-content:space-between; align-items:center;
  font-family:Consolas, monospace; font-size:20px; color:#5C6577;
  border-top:1px solid #D5DBE7; padding-top:14px;
}
.dark .footer { color:#97A0B5; border-top-color:#2C3349; }

.cols { display:flex; gap:44px; flex:1; min-height:0; }
.col { flex:1; display:flex; flex-direction:column; gap:22px; min-width:0; }

.panel {
  background:#FFFFFF; border:1px solid #D5DBE7; border-radius:4px;
  padding:30px 34px; display:flex; flex-direction:column; gap:14px;
}
.dark .panel { background:#1B2133; border-color:#2C3349; }
.panel .k {
  font-family:Consolas, monospace; font-size:21px; letter-spacing:0.12em;
  text-transform:uppercase; color:#5C6577;
}
.dark .panel .k { color:#97A0B5; }
.panel .v { font-family:Consolas, monospace; font-size:52px; color:#1A1F2E; }
.dark .panel .v { color:#E6EAF4; }
.panel .v.good { color:#2C7D57; } .dark .panel .v.good { color:#4FB98A; }
.panel .v.bad { color:#A96D12; } .dark .panel .v.bad { color:#E0A44A; }
.panel p { font-size:25px; line-height:1.45; color:#3D4658; }
.dark .panel p { color:#C3CAD9; }

ul.big { list-style:none; display:flex; flex-direction:column; gap:24px; }
ul.big li {
  font-size:31px; line-height:1.45; color:#3D4658; padding-left:38px; position:relative;
}
.dark ul.big li { color:#C3CAD9; }
ul.big li::before {
  content:""; position:absolute; left:2px; top:16px; width:14px; height:14px;
  border-radius:50%; background:#3D52C9;
}
ul.big li.good::before { background:#2C7D57; }
ul.big li.bad::before { background:#A96D12; }

.chart { background:#FFFFFF; border:1px solid #D5DBE7; border-radius:4px; padding:30px 36px; }
.dark .chart { background:#1B2133; border-color:#2C3349; }
.chart .ct { font-family:Consolas, monospace; font-size:22px; letter-spacing:0.1em; text-transform:uppercase; color:#5C6577; margin-bottom:24px; }
.dark .chart .ct { color:#97A0B5; }
.brow { display:grid; grid-template-columns:400px 1fr 150px; align-items:center; gap:24px; margin-bottom:18px; }
.brow .bn { font-size:26px; color:#3D4658; text-align:right; line-height:1.2; }
.dark .brow .bn { color:#C3CAD9; }
.brow .bt { position:relative; height:34px; background:#EDF0F7; border-radius:3px; overflow:hidden; }
.dark .brow .bt { background:#232A3D; }
.brow .bb { position:absolute; inset:0 auto 0 0; border-radius:3px; min-width:4px; }
.brow .bv { font-family:Consolas, monospace; font-size:26px; }
.b-acc { background:#3D52C9; } .b-good { background:#2C7D57; } .b-bad { background:#A96D12; } .b-grey { background:#8093C9; }

table.t { border-collapse:collapse; width:100%; background:#fff; border:1px solid #D5DBE7; border-radius:4px; overflow:hidden; }
.dark table.t { background:#1B2133; border-color:#2C3349; }
table.t th { font-family:Consolas, monospace; font-size:21px; letter-spacing:0.08em; text-transform:uppercase; color:#5C6577; background:#EDF0F7; font-weight:500; padding:16px 24px; text-align:left; }
.dark table.t th { background:#232A3D; color:#97A0B5; }
table.t td { font-size:26px; color:#3D4658; padding:15px 24px; border-top:1px solid #E4E9F2; }
.dark table.t td { color:#C3CAD9; border-top-color:#2C3349; }
table.t td.num { font-family:Consolas, monospace; }
table.t td b { color:#1A1F2E; font-weight:650; } .dark table.t td b { color:#E6EAF4; }
.good-t { color:#2C7D57 !important; font-weight:650; }
.bad-t { color:#A96D12 !important; font-weight:650; }

.figbox { flex:1; display:flex; align-items:center; justify-content:center; background:#fff; border:1px solid #D5DBE7; border-radius:4px; padding:26px; min-height:0; }
.figbox img { max-width:100%; max-height:100%; object-fit:contain; }

.quote {
  font-family:Georgia, serif; font-style:italic; font-size:40px; line-height:1.45;
  border-left:5px solid #3D52C9; padding:8px 0 8px 36px; max-width:1420px; color:#1A1F2E;
}
.dark .quote { color:#E6EAF4; border-left-color:#8FA0FF; }
.quote b { font-style:normal; font-weight:650; }
"""


def wrap(body: str, dark: bool = False, num: str = "") -> str:
    cls = "slide dark" if dark else "slide"
    return f"""<!DOCTYPE html>
<html lang="en"><head><meta charset="UTF-8"><title>SLEEP deck</title>
<meta name="hz:slide-selector" content=".slide">
<meta name="hz:canvas-width" content="1920"><meta name="hz:canvas-height" content="1080">
<style>{CSS}</style></head>
<body><div class="{cls}" data-canvas-width="1920" data-canvas-height="1080">
{body}
<div class="footer"><span>SLEEP &middot; Synaptic Learning through Error-driven Encoding &amp; Plasticity</span><span>Aditya Tripathi &middot; SP Jain &middot; {num}</span></div>
</div></body></html>"""


# ---------------------------------------------------------------------------
# Slides: (dark, body_html, speaker_notes)
# ---------------------------------------------------------------------------

SLIDES: list[tuple[bool, str, str]] = []


def slide(body: str, notes: str, dark: bool = False) -> None:
    SLIDES.append((dark, body, notes))


# 1 -- Title -----------------------------------------------------------------
slide("""
<div class="frame" style="justify-content:center;">
  <p class="eyebrow">A research journey &middot; April &ndash; August 2026</p>
  <h1 class="title" style="font-size:88px; max-width:1650px;">Recognition Without Recall:<br>
  <b>Diagnosing and Repairing</b> Biologically-Inspired Memory Consolidation in Transformers</h1>
  <p class="lead" style="font-size:38px; margin-top:26px;">Teaching large language models to remember what they read &mdash; permanently.</p>
  <p class="lead" style="font-size:30px; margin-top:60px; color:#5C6577;">Aditya Tripathi &middot; SP Jain School of Global Management<br>
  Research guidance: Prof. Debashis Guha &middot; Paper co-author: Dr. Rachit Garg</p>
</div>
""", """Good [morning/afternoon]. This is the story of a research project that set out to give language models a memory, discovered exactly why the obvious approach fails, and then repaired it. The title has two halves for a reason - the first half is a diagnosis, the second half is a fix that works. Everything I'll show you is backed by pre-registered experiments: 4 models, 3 model families, and every headline number replicated across seeds.""")

# 2 -- The problem ------------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow">01 &middot; The problem</p>
  <h1 class="title">LLMs have <b>architectural amnesia</b></h1>
  <div class="content">
    <p class="lead">A deployed language model's weights are frozen. Everything it "learns" in a conversation lives only in the context window &mdash; and dies with it.</p>
    <div class="cols" style="flex:0 0 auto;">
      <div class="panel"><span class="k">What the model has</span><p style="font-size:29px;"><b>A context window</b> &mdash; perfect short-term memory that is wiped after every conversation.</p></div>
      <div class="panel"><span class="k">What it lacks</span><p style="font-size:29px;"><b>Consolidation</b> &mdash; any path from "I just read this" to "I now know this," the way experiences become knowledge overnight for us.</p></div>
      <div class="panel"><span class="k">The usual workaround</span><p style="font-size:29px;"><b>RAG</b> bolts a database onto the prompt. Useful &mdash; but it is a lookup, not memory. The model never actually learns anything.</p></div>
    </div>
    <p class="lead">The research question: <b>can a frozen, pretrained model learn new facts permanently, from a single exposure, without destroying what it already knows?</b></p>
  </div>
</div>
""", """Set up the gap: LLMs are frozen at deployment - a form of anterograde amnesia. The context window is short-term memory that gets wiped. RAG is the industry workaround, but it's plumbing - a database lookup on every query - not learning. Our question: single-exposure, permanent, non-destructive learning in a frozen pretrained transformer. Three constraints that make it hard: one exposure only, must survive after context is gone, and must not damage existing capabilities.""")

# 3 -- The inspiration ---------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow">02 &middot; The inspiration</p>
  <h1 class="title">Biology already solved this &mdash; <b>with sleep</b></h1>
  <div class="content">
    <p class="lead">Your brain faces the same problem every day and solves it with a structured pipeline:</p>
    <div class="cols" style="flex:0 0 auto;">
      <div class="panel"><span class="k">1 &middot; Tagging</span><p>Surprising experiences get flagged at the synapse by <b>prediction-error tags</b> (synaptic tagging &amp; capture).</p></div>
      <div class="panel"><span class="k">2 &middot; Competition</span><p>Tags compete for a <b>limited pool of plasticity proteins</b> (PRPs) &mdash; only what matters gets resources.</p></div>
      <div class="panel"><span class="k">3 &middot; Replay in sleep</span><p>Winners get <b>replayed by the hippocampus</b>, which gradually trains the neocortex overnight.</p></div>
      <div class="panel"><span class="k">4 &middot; Consolidation</span><p>Episodic traces become <b>durable knowledge</b> &mdash; without overwriting old memories.</p></div>
    </div>
    <p class="lead">The result: continuous, selective, single-exposure learning that does not catastrophically forget. <b>Exactly the specification we need.</b></p>
  </div>
</div>
""", """The biology: synaptic tagging and capture (Frey & Morris) - weak stimulation creates molecular tags; tags compete for limited plasticity-related proteins; winners get consolidated. Complementary Learning Systems (McClelland): fast hippocampal store + slow neocortical store, bridged by replay during sleep. The output spec of this pipeline is exactly what we want for LLMs: continuous, selective, one-shot learning without catastrophic forgetting. So the hypothesis: implement the pipeline faithfully and the properties should follow.""")

# 4 -- The proposal ------------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow">03 &middot; The proposal</p>
  <h1 class="title">SLEEP: implement <b>all five mechanisms</b>, faithfully</h1>
  <div class="content">
    <div class="cols" style="flex:0 0 auto;">
      <div class="col">
        <ul class="big">
          <li><b>Tagging layer</b> &mdash; flags surprising spans by prediction-error z-score during normal inference ("wake")</li>
          <li><b>KV memory bank</b> &mdash; one-shot episodic store, written directly into attention</li>
          <li><b>PRP allocator</b> &mdash; a budgeted competition that decides which memories deserve consolidation</li>
        </ul>
      </div>
      <div class="col">
        <ul class="big">
          <li><b>Sleep engine</b> &mdash; replays selected memories and trains a LoRA adapter (the "neocortex")</li>
          <li><b>Safety machinery</b> &mdash; EWC + hard clipping + interleaving + validation rollback, so learning never destroys the base model</li>
        </ul>
      </div>
    </div>
    <div class="panel" style="flex:0 0 auto;"><span class="k">Discipline from day one</span><p style="font-size:29px;"><b>36 design questions resolved and pre-registered before any code was written.</b> Every hyperparameter default documented; every later deviation logged with its reason. This discipline is what made everything that follows trustworthy.</p></div>
  </div>
</div>
""", """SLEEP = Synaptic Learning through Error-driven Encoding and Plasticity. Five components, each mapping to one biological mechanism. Emphasize the pre-registration: 36 design questions resolved in a formalization document before implementation - so when things later failed, we could tell exactly which assumption broke, and when things worked, nobody could accuse us of tuning our way to the answer. About 4,000 lines of Python, 410 tests by the end.""")

# 5 -- Architecture ------------------------------------------------------------
slide(f"""
<div class="frame">
  <p class="eyebrow">04 &middot; Architecture</p>
  <h1 class="title" style="margin-bottom:26px;">One system, two phases: <b>wake and sleep</b></h1>
  <div class="figbox"><img src="{fig_uri('figure_architecture.png')}" alt="SLEEP architecture: wake phase (tagging, PRP allocation, KV writes) and sleep phase (replay, consolidation training under safety constraints)"></div>
</div>
""", """Walk the diagram left to right. WAKE: input flows through the model normally; the tagging layer watches per-token surprise; tagged spans allocate against the PRP budget and their full episode is written into the KV memory bank - this is the fast, hippocampus-like store. SLEEP: the engine replays selected memories, trains the consolidation adapter W_cons (LoRA), under four safety mechanisms - EWC penalty, hard per-parameter clipping, old-knowledge interleaving, and validation with rollback. After consolidation the bank is cleared: knowledge should now live in the weights.""")

# 6 -- First campaign ----------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow era-d">05 &middot; First campaign &middot; diagnosis era</p>
  <h1 class="title">What we did first: <b>a controlled memory exam</b></h1>
  <div class="content">
    <div class="cols" style="flex:0 0 auto;">
      <div class="panel"><span class="k">The dataset</span><p><b>200 synthetic facts</b> ("Nimbus Holdings reported Q3 revenue of..."). Fictional entities, so pretraining cannot cheat &mdash; if the model knows it, it learned it from us, once.</p></div>
      <div class="panel"><span class="k">The protocol</span><p>Single exposure during wake &rarr; full sleep cycle &rarr; test with the memory bank <b>disabled</b>. Only what reached the weights counts.</p></div>
    </div>
    <div class="cols" style="flex:0 0 auto;">
      <div class="panel"><span class="k">MC &mdash; recognition</span><p>Pick the right answer from 4 options. Chance = 0.25.</p></div>
      <div class="panel"><span class="k">DRA &mdash; recall</span><p>Answer a free-form question about the fact. The real test.</p></div>
      <div class="panel"><span class="k">BCP &mdash; damage</span><p>Perplexity after &divide; before. 1.0 = unharmed; &gt;2 = serious damage.</p></div>
    </div>
    <p class="lead">Model: Qwen2.5-7B. Three questions: does it <b>recognize</b>? Can it <b>recall</b>? Did we <b>break anything</b>?</p>
  </div>
</div>
""", """Explain the experimental design. Synthetic facts are deliberate - with real facts you can't distinguish 'learned it from us' from 'knew it from pretraining'. The three metrics matter for what comes next: MC is recognition (multiple choice), DRA is free-form recall (generation - the hard one), BCP is damage to general capability. Key protocol detail: at evaluation the KV memory bank is switched OFF - we only credit knowledge that made it into the weights.""")

# 7 -- First results -----------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow era-d">06 &middot; First results</p>
  <h1 class="title">Memory that <b>recognizes but cannot recall</b></h1>
  <div class="content">
    <div class="chart" style="flex:0 0 auto;">
      <p class="ct">Qwen2.5-7B, 200 facts, single exposure</p>
      <div class="brow"><span class="bn">Recognition &mdash; tagged facts</span><div class="bt"><div class="bb b-acc" style="width:84%"></div></div><span class="bv">0.28</span></div>
      <div class="brow"><span class="bn">Recognition &mdash; untagged</span><div class="bt"><div class="bb b-grey" style="width:36%"></div></div><span class="bv">0.12</span></div>
      <div class="brow"><span class="bn">Free-form recall (DRA)</span><div class="bt"><div class="bb b-bad" style="width:1%"></div></div><span class="bv">0.00</span></div>
    </div>
    <ul class="big">
      <li>Injected memories <b>tilt a four-way choice by +0.16</b> &mdash; encoding genuinely happens</li>
      <li>The same memories <b>never surface when the model must generate</b> the answer &mdash; recall is zero</li>
      <li class="bad">And the safety sweep found no setting with useful recall at acceptable damage &mdash; every configuration paid more damage per point of recall than naive fine-tuning</li>
    </ul>
  </div>
</div>
""", """The April headline. The recognition signal is real: +0.16 over untagged facts on multiple choice. But cloze accuracy: zero. Free-form recall: zero. The memory can bias a comparison between four options but cannot steer 50 tokens of generation against the pretrained prior. Also: the stability-plasticity sweep looked like a wall - no knee in the Pareto curve. And one more worrying observation: the system's internal validator was certifying consolidations that failed any real test.""")

# 8 -- Concluded then ----------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow era-d">07 &middot; What we concluded (then)</p>
  <h1 class="title">"Recognition without recall" &mdash; <b>but on one model, one seed</b></h1>
  <div class="content">
    <p class="quote">The transformer analogue of a classic dissociation from cognitive psychology: familiarity can be supported by a partial trace; <b>generation requires reconstructing the whole memory.</b></p>
    <ul class="big">
      <li><b>Encoding works</b> &mdash; tagging, KV injection, and resource allocation all transferred from biology</li>
      <li class="bad"><b>Consolidation doesn't</b> &mdash; nothing we stored could be recalled from the weights</li>
      <li>Open question: is this a <b>law of nature</b> for frozen models &mdash; or an artifact of how we built it?</li>
    </ul>
    <p class="lead">Everything rested on <b>one model and one seed</b>. Before believing any of it, it had to survive hostile replication.</p>
  </div>
</div>
""", """Name the finding - recognition without recall, mirroring the recognition/recall dissociation in human memory research (patients with hippocampal damage often retain recognition while losing free recall). Mechanistically: memory KV enters attention as a small perturbation - enough to flip an argmax among 4 options, overwhelmed across autoregressive generation. But be honest about the epistemic status at this point: one model, one seed. The next phase was about stress-testing every claim.""")

# 9 -- Feedback round 1 --------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow era-d">08 &middot; Hostile replication</p>
  <h1 class="title">We subjected every claim to <b>hostile replication</b></h1>
  <div class="content">
    <table class="t" style="flex:0 0 auto;">
      <tr><th>Upgrade</th><th>From &rarr; to</th><th>Why it matters</th></tr>
      <tr><td><b>Seeds</b></td><td class="num">1 &rarr; 5 per result</td><td>numbers become distributions, not anecdotes</td></tr>
      <tr><td><b>Models</b></td><td class="num">1 &rarr; 4 (3 families)</td><td>"is this just Qwen?" answered</td></tr>
      <tr><td><b>Horizon</b></td><td class="num">3 &rarr; 10 cycles</td><td>the regime the safety machinery was designed for</td></tr>
      <tr><td><b>Baselines</b></td><td>+ RAG, EWC-only, in-context</td><td>bounds the problem from above and below</td></tr>
      <tr><td><b>Validation</b></td><td>+ external recall gate</td><td>the system may no longer grade its own homework</td></tr>
    </table>
    <div class="panel" style="flex:0 0 auto;"><span class="k">Plus a ground-truth verification gate</span><p style="font-size:28px;">Before any experiment on a new architecture: three hardware-level checks (empty memory changes <b>nothing</b>, bit-identical; populated memory changes output; disable restores exactly). <b>It caught 2 real bugs before they could produce a single invalid number.</b></p></div>
  </div>
</div>
""", """This phase came out of the first review round (11 items across the faculty review). Frame it as methodology: every claim was made to survive seeds, models, horizon, and proper baselines. The two-stage validation is worth dwelling on - we added an external gate: a consolidation only counts if the model actually answers a question about the fact under greedy decoding. And the verification gate: three ground-truth checks on real hardware before any run on a new architecture; it caught a config-encoding fault and a Qwen-specific attention assumption - zero invalid numbers ever entered the record.""")

# 10 -- Extended on our own ----------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow era-d">09 &middot; Beyond the checklist</p>
  <h1 class="title">What we added <b>on our own initiative</b></h1>
  <div class="content">
    <div class="cols">
      <div class="panel"><span class="k">The warm-up extension</span><p>If the model recognizes memories but can't <b>use</b> them while generating &mdash; teach it that skill directly. A trainable memory gate + retrieval-aware training on a neutral corpus. <b>Success bar pre-registered before any run:</b> DRA 0.10&ndash;0.15 at BCP &lt; 1.5.</p></div>
      <div class="panel"><span class="k">The discovery nobody asked for</span><p>Running <b>long horizons and many seeds together</b> revealed the "plateau lottery": the safety clip always caps damage, but the level it caps at varies <b>three orders of magnitude</b> across identical configurations (2.7 &rarr; 1,531).</p></div>
    </div>
    <p class="lead" style="flex:0 0 auto;">Both were designed to be informative <b>whichever way they came out</b> &mdash; and one of them changed how we read every continual-learning paper.</p>
  </div>
</div>
""", """Two self-initiated additions. The warm-up: the repair the diagnosis implies - if the gap is a missing skill, teach the skill. We pre-registered the success criterion before running. The plateau lottery: feedback asked for more cycles and more seeds separately; doing both together revealed that the damage ceiling's LEVEL is decided by the optimization path, not the hyperparameters - identical configs landing at 2.7 or 1,531. Implication beyond our project: any continual-learning comparison reported at one horizon with one seed may be reporting the lottery, not the architecture.""")

# 11 -- Phase results ----------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow era-d">10 &middot; Replication results</p>
  <h1 class="title">The failure was <b>universal</b> &mdash; and the repair didn't repair</h1>
  <div class="content">
    <div class="chart" style="flex:0 0 auto;">
      <p class="ct">Same facts, same models: recall from context vs. recall from consolidated weights (5 seeds each)</p>
      <div class="brow"><span class="bn">Llama-8B &mdash; in context / weights</span><div class="bt"><div class="bb b-good" style="width:98.7%"></div></div><span class="bv">0.99 / 0.006</span></div>
      <div class="brow"><span class="bn">Mistral-7B &mdash; in context / weights</span><div class="bt"><div class="bb b-good" style="width:92.7%"></div></div><span class="bv">0.93 / 0.007</span></div>
      <div class="brow"><span class="bn">Qwen-7B &mdash; in context / weights</span><div class="bt"><div class="bb b-good" style="width:76%"></div></div><span class="bv">0.76 / 0.006</span></div>
      <div class="brow"><span class="bn">Qwen-1.5B &mdash; in context / weights</span><div class="bt"><div class="bb b-good" style="width:70%"></div></div><span class="bv">0.70 / 0.007</span></div>
    </div>
    <ul class="big">
      <li class="bad">Recall floor <b>&asymp; 0.006 on every model</b> &mdash; while the same models answer near-perfectly from context</li>
      <li class="bad">Internal validation <b>falsely confirms 93&ndash;100%</b> of consolidations, on every model</li>
      <li class="bad">The warm-up: <b>16 runs, 4 models &mdash; recall never moved</b> (|&Delta;DRA| &le; 0.008)</li>
    </ul>
  </div>
</div>
""", """The replication verdict. The dissociation is a property of the mechanism, not one model's pretraining: floor of 0.006 everywhere, in-context ceilings 0.70-0.99 - so it's not a comprehension failure. Proxy validation is 93-100% wrong everywhere. And the warm-up is a replicated null: 16 runs, recall never moved - though it did move recognition, which reproduces the dissociation under intervention. At this point we had a rigorous, universal diagnosis - of a failure.""")

# 12 -- Middle conclusion --------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow era-d">11 &middot; The low point</p>
  <h1 class="title">An honest negative-result paper &mdash; <b>with one uncomfortable table</b></h1>
  <div class="content">
    <table class="t" style="flex:0 0 auto;">
      <tr><th>Approach (single cycle, Qwen-7B)</th><th>Recall</th><th>Damage (BCP)</th></tr>
      <tr><td>Full SLEEP pipeline (best setting)</td><td class="num bad-t">0.012&ndash;0.103</td><td class="num">1.29&ndash;2.73</td></tr>
      <tr><td><b>Plain EWC penalty alone</b></td><td class="num">0.070</td><td class="num good-t">1.12</td></tr>
      <tr><td>Naive LoRA fine-tuning</td><td class="num">0.275</td><td class="num bad-t">2.94</td></tr>
    </table>
    <ul class="big">
      <li class="bad">A textbook regularizer &mdash; one line of math &mdash; delivered a <b>better recall-per-damage ratio than our entire five-component architecture</b></li>
      <li>The paper at this point: <i>"Recognition Without Recall: Empirical <b>Limits</b> of Biologically-Inspired Memory Consolidation"</i></li>
      <li>The question we could not shake: <b>is consolidation into frozen LLMs simply impossible &mdash; or are we doing it wrong?</b></li>
    </ul>
  </div>
</div>
""", """The low point, stated plainly. We had found the uncomfortable comparison before a reviewer did: EWC alone beat our whole pipeline on recall-per-damage. The paper title at this stage said 'Empirical Limits'. A defensible, rigorous negative result - but the nagging question: law of nature, or wrong implementation? That tension sets up the turn.""")

# 13 -- The reframe --------------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow era-r">12 &middot; The turn</p>
  <h1 class="title">A pre-mortem: <b>"assume the paper is wrong &mdash; where is the error?"</b></h1>
  <div class="content">
    <p class="lead" style="flex:0 0 auto;">Before closing the project, we re-examined every design choice from scratch, with fresh eyes. One assumption had never been questioned: <b>where</b> we were writing.</p>
    <div class="cols" style="flex:0 0 auto;">
      <div class="panel"><span class="k">The unexamined choice</span><p>We placed the consolidation adapter on <b>top-third attention layers</b> &mdash; because that's the "neocortex" in the biological analogy.</p></div>
      <div class="panel"><span class="k">What the evidence says</span><p>Causal tracing (ROME, MEMIT, Geva et&nbsp;al.) locates factual associations in <b>mid-stack MLP layers</b>. We were writing to the wrong address.</p></div>
    </div>
    <ul class="big" style="flex:0 0 auto;">
      <li class="good"><b>Move the write</b> &mdash; consolidation adapter to mid-stack MLPs, where facts actually live</li>
      <li class="good"><b>Teach the fact, not the string</b> &mdash; train on 20+ paraphrases including question&ndash;answer forms</li>
      <li class="good"><b>Distill from the model's own best self</b> &mdash; the in-context teacher already answers at 0.76&ndash;0.99; make the adapter match it</li>
    </ul>
  </div>
</div>
""", """The reframe. Present it as a pre-mortem exercise: assume failure was our fault - where's the bug? The one never-questioned choice was the biological analogy for PLACEMENT: top-third attention as 'neocortex'. The knowledge-editing literature says facts live in mid-stack MLPs. Then two more first-principles corrections: single-wording training teaches string reproduction, not knowledge (we saw cloze succeed while recall failed - the signature of memorization); and the model already 'knows' every fact when it can see it - so distill from the in-context teacher. Plus one hygiene fix: greedy decoding for all knowledge claims.""")

# 14 -- Ladder --------------------------------------------------------------------
slide(f"""
<div class="frame">
  <p class="eyebrow era-r">13 &middot; What we saw</p>
  <h1 class="title" style="margin-bottom:24px;">The ladder: <b>0.01 &rarr; 0.36</b>, each factor isolated</h1>
  <div class="figbox" style="margin-bottom:22px;"><img src="{fig_uri('figure_localisation_ladder.png')}" alt="Five-arm ladder: recall rises and damage falls as placement, paraphrases, and distillation are added"></div>
  <p class="lead" style="flex:0 0 auto; font-size:29px;">Placement halves damage in <b>8 of 8</b> comparisons &middot; paraphrases flip cloze <b>down</b> and recall <b>up</b> (memorization &rarr; extraction) &middot; distillation lands at <b>0.36 recall with BCP 0.83</b> &mdash; better than baseline preservation. The control shows the objective cannot rescue the wrong substrate.</p>
</div>
""", """The isolation ladder - five arms, two models, two seeds, matched training budget. Read it left to right: original placement 0.01; move to mid-MLP alone and damage halves (8/8 comparisons); add paraphrases and recall jumps to 0.22 while CLOZE falls - that opposite movement is the smoking gun of memorization becoming extraction; add distillation: 0.36 recall at BCP 0.83 - better preservation than baseline. And the control arm - distillation on the OLD placement - recovers almost nothing: the objective can't rescue the wrong substrate. Sixty times the old floor, each factor's contribution isolated.""")

# 15 -- Third cause ---------------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow era-r">14 &middot; What we saw, continued</p>
  <h1 class="title">A third culprit: <b>our own safety machinery</b></h1>
  <div class="content">
    <p class="lead" style="flex:0 0 auto;">Wired into the full autonomous pipeline, the winning recipe... collapsed back to 0.007. Instead of accepting it, we A/B-tested the machinery itself:</p>
    <div class="chart" style="flex:0 0 auto;">
      <p class="ct">Identical pipeline, 3 seeds each &mdash; only the safety calibration differs (trained-subset DRA)</p>
      <div class="brow"><span class="bn">Safety as designed (&delta;=0.01)</span><div class="bt"><div class="bb b-bad" style="width:2.4%"></div></div><span class="bv">0.007</span></div>
      <div class="brow"><span class="bn">Safety relaxed</span><div class="bt"><div class="bb b-acc" style="width:73%"></div></div><span class="bv">0.210 &plusmn; 0.018</span></div>
      <div class="brow"><span class="bn">Safety moderate (&delta;=0.1)</span><div class="bt"><div class="bb b-good" style="width:66%"></div></div><span class="bv">0.189 &plusmn; 0.031 @ BCP 0.88</span></div>
    </div>
    <ul class="big">
      <li>The designed clip was so tight it <b>erased the write after training performed it</b></li>
      <li class="good">A moderate clip keeps <b>both</b> the recall and the protection &mdash; every seed under BCP 1.0</li>
      <li>The original failure had <b>three stacked causes</b>: wrong substrate, wrong wording, over-tight clip. <b>Any one alone keeps recall at the floor</b> &mdash; which is why no single fix ever moved it.</li>
    </ul>
  </div>
</div>
""", """The third cause. Integration initially erased everything - back to 0.007. The killer detail: the hard clip (delta_max 0.01) was clipping away the consolidation write AFTER training performed it. Relaxed: 0.210. Moderate (0.1): 0.189 at BCP 0.88 - every seed below 1.0, first externally-verified autonomous consolidations (21-24 of ~50 facts passing the direct recall gate per run). Big picture: three stacked causes - substrate, wording, clip. Each alone suffices to floor recall. That's why the first movement's single-factor interventions all failed - and why the diagnosis took a full pre-mortem to crack.""")

# 16 -- Matrix ---------------------------------------------------------------------
slide(f"""
<div class="frame">
  <p class="eyebrow era-r">15 &middot; The proof</p>
  <h1 class="title" style="margin-bottom:24px;">Pre-registered validation: <b>every model, 17&ndash;125&times; the old floor</b></h1>
  <div class="figbox" style="margin-bottom:22px;"><img src="{fig_uri('figure_arc.png')}" alt="Before/after: universal 0.006 floor versus repaired recall of 0.105 to 0.750 across four models"></div>
  <p class="lead" style="flex:0 0 auto; font-size:29px;">35 runs, gates and predictions committed <b>before launch</b>, zero failures. Best result: <b>Mistral-7B at 0.750 &plusmn; 0.033</b> &mdash; a family the recipe never saw during development. Llama honestly flagged: recall 0.512 but damage above our bar &mdash; <b>"not yet tuned," in the paper's own words.</b></p>
</div>
""", """The validation matrix - the whole campaign was pre-registered: gates, predictions, stop-loss committed to the repository before launch; 35 runs, zero failures. Left panel: the old story. Right: the repaired system - 0.750 on Mistral (5 seeds, plus-minus 0.033), 0.512 Llama, 0.189 Qwen-7B, 0.105 at 1.5B. Two honest notes: strongest result is on a family the recipe never touched in development (argues against overfitting the recipe to Qwen); and Llama fails our damage bar at every clip setting we pre-registered - it's flagged 'not yet tuned', not hidden. The 1.5B result suggests extraction capacity grows with scale.""")

# 17 -- Ten cycles + EWC ------------------------------------------------------------
slide(f"""
<div class="frame">
  <p class="eyebrow era-r">16 &middot; The proof, long-run</p>
  <h1 class="title" style="margin-bottom:24px;">Ten cycles: <b>SLEEP beats both baselines</b> in its design regime</h1>
  <div class="figbox" style="margin-bottom:22px;"><img src="{fig_uri('figure_matrix_multi_cycle.png')}" alt="Ten-cycle trajectories: SLEEP v2 holds ~3x cumulative recall above EWC-only and naive LoRA on both seeds"></div>
  <p class="lead" style="flex:0 0 auto; font-size:29px;">Beats naive fine-tuning on cumulative recall on <b>4 of 4 models, both seeds</b> &middot; beats EWC-only <b>3&times;</b> (0.174 vs 0.056) &mdash; <b>regularization prevents forgetting; it does not create knowledge</b> &middot; and at 1.5B, naive training detonates (damage &times;24.5) while SLEEP holds &times;2.4.</p>
</div>
""", """The long-run proof, plotted from the released run files. SLEEP v2 holds ~3x cumulative recall above BOTH baselines for all ten cycles, on both seeds. The EWC comparison is the anti-strawman result: matched arms - identical adapter, steps, learning rate, only the penalty differs. EWC gets the best damage numbers but recall stays at naive's level: preservation and acquisition are different capabilities. Bonus findings: at 1.5B naive LoRA blows up (BCP 24.5 / 11.4 - the plateau lottery reproduced) while SLEEP holds ~2.4 - the safety machinery is load-bearing exactly where models are fragile. And SLEEP's across-seed spread is ~0.005 everywhere - the seed lottery turns out to belong to UNCONSTRAINED arms.""")

# 18 -- What changed ------------------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow era-r">17 &middot; What all changed</p>
  <h1 class="title">Before &rarr; after</h1>
  <div class="content">
    <table class="t">
      <tr><th style="width:44%">Before (the diagnosis era)</th><th>After (the repaired system)</th></tr>
      <tr><td>Recall floor <b>0.006</b>, universal</td><td class="good-t">0.105 &ndash; 0.750 across 4 models, 5 seeds each</td></tr>
      <tr><td>Safety machinery erased learning</td><td class="good-t">Right-sized, it prevents a &times;24 damage catastrophe</td></tr>
      <tr><td>Damage plateau: a seed lottery (2.7 &rarr; 1,531)</td><td class="good-t">Seed spread &asymp; 0.005 &mdash; the lottery belongs to unconstrained baselines</td></tr>
      <tr><td>Validator 93&ndash;100% wrong</td><td class="good-t">External greedy recall gate &mdash; every consolidation independently verified</td></tr>
      <tr><td>Title: "...Empirical <b>Limits</b>..."</td><td class="good-t">Title: "...<b>Diagnosing and Repairing</b>..."</td></tr>
    </table>
    <p class="quote" style="flex:0 0 auto; font-size:36px;">The meta-lesson: biological <b>mechanisms</b> transfer &mdash; tagging, selective replay, staged consolidation. Biological <b>anatomy</b> does not. The transformer has its own anatomy, and the editing literature had already mapped it.</p>
  </div>
</div>
""", """The before/after summary - including the paper's own title changing from 'Empirical Limits' to 'Diagnosing and Repairing'. Land the meta-lesson hard, because it generalizes beyond our project: when transferring biology to ML, copy the mechanisms (tagging, replay, staged consolidation, sleep scheduling) but let the substrate's own causal structure decide the wiring. The hippocampus-to-top-layers analogy was the single most expensive assumption in the project.""")

# 19 -- Final conclusion ---------------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow era-r">18 &middot; Conclusion</p>
  <h1 class="title">SLEEP works in its design regime &mdash; <b>with named limits</b></h1>
  <div class="content">
    <ul class="big">
      <li class="good">A frozen LLM can now <b>read once, sleep, and answer questions from its own weights</b> &mdash; 0.75 recall on its best family, verified consolidations every cycle</li>
      <li class="good">Beats naive fine-tuning (4/4 models) <b>and</b> the standard continual-learning baseline (3&times; vs EWC) under pre-registered gates</li>
      <li class="good">The diagnostic findings stand on their own: recognition-without-recall, the proxy-validation failure, the seed-instability of unconstrained baselines</li>
    </ul>
    <div class="cols" style="flex:0 0 auto;">
      <div class="panel"><span class="k">Open &middot; 1</span><p><b>Early memories fade</b> &mdash; the engine rehearses only the current batch. Interleaved rehearsal (what biological sleep actually does) is the clearest next mechanism.</p></div>
      <div class="panel"><span class="k">Open &middot; 2</span><p><b>Damage still leaks</b> &asymp; +0.12/cycle in every arm &mdash; no configuration holds BCP &lt; 1.05 over ten cycles.</p></div>
      <div class="panel"><span class="k">Open &middot; 3</span><p><b>Llama needs per-family tuning</b> &mdash; its bleed is not clip-dominated; learning-rate work required.</p></div>
    </div>
  </div>
</div>
""", """The final conclusion, calibrated. What we claim: the system works in its design regime, validated across families and scales, against real baselines, under pre-registered gates. What we don't claim: solved. Three named gaps, each with a specific next experiment: interleaved rehearsal for the early-batch fade (this is also the most biologically faithful next step - real sleep interleaves old memories with new), the damage leak, and Llama's per-family tuning. Naming your own gaps is what makes the positive claims believable.""")

# 20 -- Paper -----------------------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow">19 &middot; The paper</p>
  <h1 class="title">Written, revised, and <b>ready for review</b></h1>
  <div class="content">
    <div class="cols" style="flex:0 0 auto;">
      <div class="panel"><span class="k">The paper</span><p style="font-size:28px;"><b>"Recognition Without Recall: Diagnosing and Repairing Biologically-Inspired Memory Consolidation in Transformers"</b><br><br>31 pages &middot; two movements (diagnosis &rarr; repair) &middot; 9 figures &middot; every claim with seeds and its number</p></div>
      <div class="col">
        <div class="panel"><span class="k">Evidence base</span><p style="font-size:28px;">4 models &times; 3 families &middot; 95+ seeded runs &middot; pre-registered gates &middot; ~$43 total GPU spend</p></div>
        <div class="panel"><span class="k">Fully public</span><p style="font-size:28px;">Implementation (410 tests), formalization, pre-registrations, all result files &amp; logs &mdash; every number in the paper traces to a file in the repository</p></div>
      </div>
    </div>
    <p class="lead" style="flex:0 0 auto;">The six formalization amendments &mdash; including the localisation correction itself &mdash; are documented as first-class results: <b>the record of what empirical work forced us to unlearn.</b></p>
  </div>
</div>
""", """The paper status: 31 pages, restructured as two movements so the reader experiences the arc in order - diagnosis with full rigor, then repair with pre-registered validation. Reproducibility is total: the repo has the implementation with 410 tests, the 36-question formalization, the pre-registration documents, and every result JSON and log - any number in the paper can be traced to a file. Entire empirical programme cost about $43 of cloud GPU - rigor at this scale is accessible.""")

# 21 -- Future ------------------------------------------------------------------------
slide("""
<div class="frame">
  <p class="eyebrow">20 &middot; Why it matters</p>
  <h1 class="title">Where this could go</h1>
  <div class="content">
    <div class="cols" style="flex:0 0 auto;">
      <div class="panel"><span class="k">Assistants that know you</span><p>A model that consolidates what you told it <b>last month</b> &mdash; no prompt stuffing, no retrieval index, no re-explaining yourself.</p></div>
      <div class="panel"><span class="k">"Sleep on it" for enterprises</span><p>Models that ingest the day's documents <b>overnight</b> and answer from weights tomorrow &mdash; learning as an operating rhythm, not a retraining project.</p></div>
      <div class="panel"><span class="k">On-device learning</span><p>Small models are exactly where naive fine-tuning detonates (&times;24 damage) and where right-sized safety machinery <b>earned its keep</b>.</p></div>
    </div>
    <div class="cols" style="flex:0 0 auto;">
      <div class="panel"><span class="k">The safety half</span><p>Consolidation that works means <b>silent data retention</b> is now a real concern: deployed systems need user control over what is encoded, consolidated, and kept &mdash; plus external validation gates, which our results show are not optional.</p></div>
      <div class="panel"><span class="k">Next experiment</span><p><b>Interleaved rehearsal</b> &mdash; replay old consolidations alongside new ones during sleep. The most biologically faithful mechanism we haven't yet built, aimed at the one gap that matters most.</p></div>
    </div>
  </div>
</div>
""", """Real-world implications, kept honest. Personal assistants with genuine long-term memory; enterprise 'sleep on it' ingestion - learning as a nightly rhythm instead of a retraining project; on-device continual learning, where our small-model result is most relevant. Give the safety paragraph its full weight: a working consolidation path makes silent data retention a real concern - user controls and external validation gates should ship with the capability. And close the research loop: interleaved rehearsal is next.""")

# 22 -- Thank you ------------------------------------------------------------------------
slide("""
<div class="frame" style="justify-content:center;">
  <p class="eyebrow">Thank you</p>
  <h1 class="title" style="font-size:80px;">From <b>"it cannot recall"</b><br>to <b>"move the write, and it can."</b></h1>
  <p class="lead" style="font-size:32px; margin-top:40px;">Code, paper, pre-registrations &amp; all results:<br><b>github.com/Adineu03/sleep-framework</b></p>
  <p class="lead" style="font-size:28px; margin-top:50px; color:#5C6577;">Aditya Tripathi &middot; adityatripathi1503@gmail.com<br>SP Jain School of Global Management</p>
</div>
""", """Close with the one-line arc: from 'it cannot recall' to 'move the write, and it can.' Point to the repository - everything is public and reproducible. Thank the audience; thank Prof. Guha for research guidance and Dr. Garg for collaboration on the paper. Questions welcome - likely ones: why synthetic facts (pretraining contamination), why Llama fails the damage bar (per-family lr, named gap), how big is the deck's headline number's evidence (5 seeds, pre-registered).""", )


# ---------------------------------------------------------------------------
# Render + assemble
# ---------------------------------------------------------------------------

def render() -> list[str]:
    pngs = []
    for i, (dark, body, _notes) in enumerate(SLIDES, 1):
        html_path = os.path.join(HTML_DIR, f"slide_{i:02d}.html")
        png_path = os.path.join(PNG_DIR, f"slide_{i:02d}.png")
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(wrap(body, dark=dark, num=f"{i:02d} / {len(SLIDES)}"))
        url = "file:///" + html_path.replace("\\", "/")
        subprocess.run([
            EDGE, "--headless", "--disable-gpu", "--hide-scrollbars",
            "--force-device-scale-factor=2",
            f"--screenshot={png_path}", "--window-size=1920,1080",
            "--default-background-color=00000000", url,
        ], check=True, capture_output=True, timeout=60)
        if not os.path.exists(png_path):
            raise RuntimeError(f"render failed for slide {i}")
        pngs.append(png_path)
        print(f"rendered slide {i:02d}")
    return pngs


def assemble(pngs: list[str]) -> str:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for png, (_dark, _body, notes) in zip(pngs, SLIDES):
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
        s.notes_slide.notes_text_frame.text = notes

    out = os.path.join(HERE, "SLEEP_Presentation.pptx")
    prs.save(out)
    return out


if __name__ == "__main__":
    files = render()
    out = assemble(files)
    print(f"\nDeck written: {out} ({len(files)} slides)")
